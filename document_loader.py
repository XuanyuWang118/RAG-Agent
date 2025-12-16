import os
import io
from typing import List, Dict, Optional

import fitz
import pdfplumber
import docx2txt
from pptx import Presentation
from PIL import Image

from config import DATA_DIR


class DocumentLoader:
    def __init__(
        self,
        data_dir: str = DATA_DIR,
        # 新增一个参数用于存放提取出的图片
        image_output_dir: str = "./images_extracted", 
    ):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]
        self.image_output_dir = image_output_dir
        os.makedirs(self.image_output_dir, exist_ok=True) # 确保图片输出目录存在

    def load_pdf(self, file_path: str) -> List[Dict]:
        """
        加载PDF文件，按页返回内容，并使用 PyMuPDF 提取图片和保存到本地。
        最终返回的字典只包含 'text' 和 'images' 键，以供 load_document 统一封装元数据。
        """
        pages = []
        filename = os.path.basename(file_path)
        filename_base = os.path.splitext(filename)[0]

        # 封装图片保存逻辑作为内部辅助函数
        def _save_image_from_bytes(image_bytes: bytes, page_num: int, xref: int, ext: str) -> str:
            """将图片二进制数据保存到本地"""
            if not image_bytes:
                return None
            
            ext = ext.lower() if ext.lower() in ['png', 'jpg', 'jpeg', 'bmp', 'tiff'] else 'png'
            
            # 使用 PyMuPDF 的 xref 作为标识符
            img_filename = f"{filename_base}_p{page_num}_xref_{xref}.{ext}"
            img_save_path = os.path.join(self.image_output_dir, img_filename)
            
            try:
                # 使用 PIL 从字节流中加载并保存图片
                img = Image.open(io.BytesIO(image_bytes))
                img.save(img_save_path)
                return img_save_path
            except Exception:
                # 提取或保存失败
                return None

        try:
            # 1. 使用 PyMuPDF (fitz) 打开 PDF 文件
            with fitz.open(file_path) as pdf_document:
                for i, page in enumerate(pdf_document):
                    image_info = []
                    page_num = i + 1
                    
                    # 2. 提取文本内容
                    text = page.get_text() or ""
                    
                    # 3. 提取和保存图片
                    for img_tuple in page.get_images(full=True):
                        xref = img_tuple[0]
                        img_ext = img_tuple[7]
                        
                        try:
                            image_data = pdf_document.extract_image(xref)
                            if image_data and image_data['image']:
                                image_bytes = image_data['image']
                                
                                # 调用内部函数保存图片
                                img_save_path = _save_image_from_bytes(
                                    image_bytes, 
                                    page_num, 
                                    xref, 
                                    img_ext
                                )
                                
                                if img_save_path:
                                    # 记录图片路径信息
                                    # 注意：原 load_pdf 结构中没有 'name' 字段，但为了可追溯性保留
                                    image_info.append({"path": img_save_path, "name": f"xref_{xref}"})
                                    
                        except Exception:
                            continue

                    # 4. 格式化文本
                    formatted_text = f"--- 第 {page_num} 页 ---\n{text.strip()}\n"
                    
                    # 5. 存储结果：只返回 load_document 需要的核心数据
                    pages.append({
                        "text": formatted_text,
                        "images": image_info  
                    })
                    
        except Exception as e:
            print(f"加载PDF文件失败 {file_path}: {e}")
            return []
            
        return pages

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，按幻灯片返回内容

        TODO: 实现PPT文件加载
        要求：
        1. 使用Presentation读取PPT文件
        2. 遍历每一页，提取文本内容
        3. 格式化为"--- 幻灯片 X ---\n文本内容\n"
        4. 返回幻灯片内容列表，每个元素包含 {"text": "..."}
        """
        slides_content = []
        try:
            # 1. 使用Presentation读取PPT文件
            prs = Presentation(file_path)
            filename_base = os.path.basename(file_path).replace(os.path.splitext(file_path)[1], "")
            
            # 2. 遍历每一页，提取文本内容
            for i, slide in enumerate(prs.slides):
                slide_text = []
                image_info = []
                
                for shape in slide.shapes:
                    # 提取文本内容
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            # 拼接所有文本块内容
                            slide_text.append("".join(run.text for run in paragraph.runs))
                    
                    # 提取图片内容
                    if shape.shape_type == 13: # 形状类型 13 代表 PICTURE
                        try:
                            # 检查形状是否真的包含图片
                            if hasattr(shape, "image"):
                                image_part = shape.image
                                image_bytes = image_part.blob
                                image_ext = image_part.ext

                                # 保存图片
                                img_filename = f"{filename_base}_s{i+1}_{shape.name}.{image_ext}"
                                img_save_path = os.path.join(self.image_output_dir, img_filename)

                                with open(img_save_path, 'wb') as f:
                                    f.write(image_bytes)

                                image_info.append({"path": img_save_path, "name": shape.name})
                        except (AttributeError, Exception) as img_e:
                            # 有些形状被误识别为图片，跳过这些错误
                            continue

                text = "\n".join(filter(None, slide_text)).strip()
                formatted_text = f"--- 幻灯片 {i + 1} ---\n{text}\n"
                
                slides_content.append({
                    "text": formatted_text,
                    "images": image_info # 新增图片信息字段
                })
        except Exception as e:
            print(f"加载PPTX文件失败 {file_path}: {e}")
        return slides_content

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件
        TODO: 实现DOCX文件加载
        要求：
        1. 使用docx2txt读取DOCX文件
        2. 返回文本内容
        """
        try:
            # 1. 使用docx2txt读取DOCX文件并返回文本内容
            text = docx2txt.process(file_path)
            # 2. 返回文本内容
            return text.strip()
        except Exception as e:
            print(f"加载DOCX文件失败 {file_path}: {e}")
            return ""

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件
        TODO: 实现TXT文件加载
        要求：
        1. 使用open读取TXT文件（注意使用encoding="utf-8"）
        2. 返回文本内容
        """
        try:
            # 1. 使用open读取TXT文件（注意使用encoding="utf-8"）
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            # 2. 返回文本内容
            return text.strip()
        except Exception as e:
            print(f"加载TXT文件失败 {file_path}: {e}")
            return ""

    def load_document(self, file_path: str) -> List[Dict[str, any]]:
        """加载单个文档，PDF和PPT按页/幻灯片分割，返回文档块列表"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append(
                    {
                        "content": page_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": page_idx,
                        "images": page_data["images"], # 传入图片信息
                    }
                )
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append(
                    {
                        "content": slide_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": slide_idx,
                        "images": slide_data["images"], # 传入图片信息
                    }
                )
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                    }
                )
        else:
            print(f"不支持的文件格式: {ext}")

        return documents

    def load_all_documents(self) -> List[Dict[str, any]]:
        """加载数据目录下的所有文档"""
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return None

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        return documents

    def process_uploaded_file(self, uploaded_file) -> List[Dict]:
        """处理单个上传的文件（来自Streamlit），复用现有的load方法

        参数:
            uploaded_file: Streamlit上传的文件对象

        返回:
            处理后的文档块列表
        """
        try:
            # 获取文件扩展名
            file_name = uploaded_file.name
            file_ext = os.path.splitext(file_name)[1].lower()

            if file_ext not in self.supported_formats:
                raise ValueError(f"不支持的文件格式: {file_ext}")

            # 创建临时文件供现有方法使用
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                tmp_file_path = tmp_file.name

            try:
                # 复用现有的load方法
                if file_ext == ".pdf":
                    pages = self.load_pdf(tmp_file_path)
                    # 转换格式以匹配期望的输出
                    converted_pages = []
                    for page in pages:
                        converted_pages.append({
                            "content": page["text"],
                            "filename": file_name,
                            "filepath": f"uploaded://{file_name}",
                            "filetype": ".pdf",
                            "page_number": page.get("page_number", 0),
                            "images": page.get("images", [])
                        })
                    return converted_pages

                elif file_ext == ".pptx":
                    slides = self.load_pptx(tmp_file_path)
                    # 转换格式
                    converted_slides = []
                    for slide in slides:
                        converted_slides.append({
                            "content": slide["text"],
                            "filename": file_name,
                            "filepath": f"uploaded://{file_name}",
                            "filetype": ".pptx",
                            "page_number": slide.get("page_number", 0),
                            "images": slide.get("images", [])
                        })
                    return converted_slides

                elif file_ext == ".docx":
                    # DOCX文件直接返回文本内容
                    text = self.load_docx(tmp_file_path)
                    return [{"content": text, "filename": file_name, "filetype": file_ext, "page_number": 0}]

                elif file_ext == ".txt":
                    # TXT文件直接返回文本内容
                    text = self.load_txt(tmp_file_path)
                    return [{"content": text, "filename": file_name, "filetype": file_ext, "page_number": 0}]

                else:
                    raise ValueError(f"不支持的文件类型: {file_ext}")

            finally:
                # 清理临时文件
                os.unlink(tmp_file_path)

        except Exception as e:
            print(f"处理上传文件失败 {uploaded_file.name}: {e}")
            return []
